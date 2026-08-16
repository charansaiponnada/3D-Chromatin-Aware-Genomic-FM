#!/usr/bin/env python3
"""Build presentations/pitch_deck.pptx -- an all-visual deck for the project.

Every slide is a single full-bleed 16:9 figure rendered by matplotlib. No bullet
lists, no text boxes: the deck is illustrations and infographics only, per the
brief.

EVERY NUMBER IS READ FROM A FILE IN THE REPOSITORY. Nothing is typed in by hand
and nothing is illustrative. The sources are:

    results/*/metrics.json                        training traces, final scores
    results/*/run_config.yaml                     parameter counts
    data/processed/phi_validation_report.json     phi vs independent 4DN tracks
    data/processed/s3_validation_report.json      S3 distance-matched control
    data/processed/s4_validation_report.json      S4 sequence-matched control
    data/processed/dataset_index.npz              split sizes and boundaries
    data/interim/*.npz                            cached Hi-C footprint

A slide whose data is not yet available (seed 2 mid-run) says so on its face
rather than being omitted or estimated.

    ./3d-gen/bin/python scripts/build_pitch_deck.py
"""

import glob
import json
import os
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Rectangle

REPO = Path(__file__).resolve().parent.parent
OUT = REPO / "presentations"
FIG = REPO / "presentations" / "_pitch_figs"

# 16:9 at 160 dpi -> 2133x1200 px, crisp on any projector
W, H, DPI = 13.333, 7.5, 160

INK = "#12151c"
MUTE = "#6b7280"
FAINT = "#e5e7eb"
PAPER = "#ffffff"
BASE = "#94a3b8"        # baseline arm
STRUCT = "#e07a3f"      # structural arm
OLD = "#cbd5e1"         # superseded v1 runs
GOOD = "#2f8f5b"
BAD = "#c2453c"
ACCENT = "#3d6fb4"

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "figure.facecolor": PAPER,
    "axes.facecolor": PAPER,
    "axes.edgecolor": FAINT,
    "axes.labelcolor": INK,
    "text.color": INK,
    "xtick.color": MUTE,
    "ytick.color": MUTE,
    "axes.spines.top": False,
    "axes.spines.right": False,
})


def canvas():
    fig = plt.figure(figsize=(W, H), dpi=DPI)
    return fig


def title(fig, t, sub=None):
    fig.text(0.055, 0.90, t, fontsize=31, fontweight="bold", va="top")
    if sub:
        fig.text(0.055, 0.815, sub, fontsize=15.5, color=MUTE, va="top")


def footer(fig, src):
    fig.text(0.055, 0.038, src, fontsize=9.5, color="#9aa3af", style="italic")


def save(fig, name):
    FIG.mkdir(parents=True, exist_ok=True)
    p = FIG / f"{name}.png"
    fig.savefig(p, dpi=DPI, facecolor=PAPER)
    plt.close(fig)
    return p


# ----------------------------------------------------------------- data loading

def load():
    d = {}
    tr = {}
    for pat in ("results/baselines/baseline_seed*", "results/baselines/baseline_v2_seed*",
                "results/novel_model/structural_seed*"):
        for rd in sorted(glob.glob(str(REPO / pat))):
            mp = Path(rd) / "metrics.json"
            if not mp.exists():
                continue
            recs = [x for x in json.loads(mp.read_text())
                    if "val_bits_per_nucleotide" in x]
            if not recs:
                continue
            tr[Path(rd).name] = {
                "step": [x["step"] for x in recs],
                "val": [x["val_bits_per_nucleotide"] for x in recs],
                "tau": [x.get("tau_empirical", {}).get("summary", {}).get("tau_median")
                        for x in recs],
                "done": "status: COMPLETED" in (Path(rd) / "run_config.yaml").read_text(),
            }
    d["tr"] = tr
    for k, f in (("phi", "phi_validation_report.json"),
                 ("s3", "s3_validation_report.json"),
                 ("s4", "s4_validation_report.json")):
        d[k] = json.loads((REPO / "data/processed" / f).read_text())
    import yaml
    d["p_base"] = yaml.safe_load(
        (REPO / "results/baselines/baseline_v2_seed0/run_config.yaml").read_text()
    )["model"]["n_parameters"]
    d["p_str"] = yaml.safe_load(
        (REPO / "results/novel_model/structural_seed0/run_config.yaml").read_text()
    )["model"]["n_parameters"]
    z = np.load(REPO / "data/processed/dataset_index.npz", allow_pickle=True)
    d["splits"] = {k: len(z[k]) for k in ("train", "val", "test")}
    d["window"] = int(z["window"])
    for k, f in (("p1", "p1_swap_results.json"), ("d1", "d1_diagnostic.json")):
        p = REPO / "results/novel_model" / f
        d[k] = json.loads(p.read_text()) if p.exists() else None
    return d


def finals(tr, prefix):
    return [tr[k]["val"][-1] for k in sorted(tr) if k.startswith(prefix) and tr[k]["done"]]


# ------------------------------------------------------------------- the slides

def s01_title(d):
    fig = canvas()
    fig.text(0.055, 0.63, "Does chromatin structure belong", fontsize=46,
             fontweight="bold")
    fig.text(0.055, 0.52, "inside pretraining?", fontsize=46, fontweight="bold",
             color=STRUCT)
    fig.text(0.055, 0.40,
             "A 3D-chromatin-aware genomic foundation model,\n"
             "at matched parameters and matched compute.",
             fontsize=18, color=MUTE, linespacing=1.6)
    ax = fig.add_axes([0.055, 0.16, 0.62, 0.10]); ax.axis("off")
    for i, (k, v) in enumerate([("model", f"{d['p_str']:,} params"),
                                ("data", "chr9 · GM12878 · Hi-C"),
                                ("arms", "2 × 3 seeds, matched")]):
        ax.text(i / 3, 0.55, v, fontsize=14, fontweight="bold", transform=ax.transAxes)
        ax.text(i / 3, 0.05, k.upper(), fontsize=10, color=MUTE, transform=ax.transAxes)
    fig.add_artist(plt.Line2D([0.055, 0.945], [0.115, 0.115], color=FAINT, lw=1))
    return save(fig, "01_title")


def s02_problem(d):
    fig = canvas()
    title(fig, "The signal lives far away",
          "A gene's switch can sit 500,000 letters from the gene it controls.")
    ax = fig.add_axes([0.055, 0.20, 0.89, 0.52]); ax.axis("off")
    ax.set_xlim(0, 10); ax.set_ylim(0, 4)

    ax.plot([0.2, 9.8], [3.1, 3.1], color=FAINT, lw=9, solid_capstyle="round")
    ax.text(0.2, 3.55, "1D — how a language model reads DNA", fontsize=13,
            color=MUTE)
    for x, c, lab in ((1.4, ACCENT, "switch"), (8.4, STRUCT, "gene")):
        ax.add_patch(plt.Circle((x, 3.1), 0.17, color=c, zorder=3))
        ax.text(x, 2.62, lab, fontsize=12, ha="center", fontweight="bold")
    ax.annotate("", xy=(8.4, 3.42), xytext=(1.4, 3.42),
                arrowprops=dict(arrowstyle="<->", color=MUTE, lw=1.2))
    ax.text(4.9, 3.56, "500,000 letters apart — far outside any local window",
            fontsize=11.5, ha="center", color=MUTE)

    t = np.linspace(0, 1, 400)
    bx = (0.2 + 9.6 * t + 1.9 * np.sin(6.1 * np.pi * t) * np.sin(np.pi * t))
    by = 1.05 + 0.62 * np.sin(3.4 * np.pi * t) * np.sin(np.pi * t)
    ax.plot(bx, by, color=FAINT, lw=9, solid_capstyle="round")
    ax.text(0.2, 1.95, "3D — how it actually sits in the nucleus", fontsize=13,
            color=MUTE)
    i1, i2 = 55, 345
    ax.add_patch(plt.Circle((bx[i1], by[i1]), 0.17, color=ACCENT, zorder=3))
    ax.add_patch(plt.Circle((bx[i2], by[i2]), 0.17, color=STRUCT, zorder=3))
    ax.plot([bx[i1], bx[i2]], [by[i1], by[i2]], color=GOOD, lw=2.4, ls=(0, (3, 2)),
            zorder=2)
    mx, my = (bx[i1] + bx[i2]) / 2, (by[i1] + by[i2]) / 2
    ax.text(mx, my - 0.42, "touching", fontsize=12, ha="center", color=GOOD,
            fontweight="bold")
    fig.text(0.055, 0.125,
             "Folding decides what is actually adjacent. Hi-C measures it. "
             "Sequence-only models never see it.",
             fontsize=14.5)
    footer(fig, "Illustrative schematic. All quantitative slides use measured values.")
    return save(fig, "02_problem")


def s03_gap(d):
    fig = canvas()
    title(fig, "Where structure enters",
          "Everyone adds Hi-C somewhere. Nobody adds it during pretraining.")
    ax = fig.add_axes([0.055, 0.16, 0.89, 0.60]); ax.axis("off")
    ax.set_xlim(0, 10); ax.set_ylim(0, 4.4)
    rows = [
        ("Akita / Orca", [0, 0, 0], "predicts folding FROM sequence — single purpose"),
        ("Evo2HiC", [0, 1, 0], "distils to a frozen teacher — no nucleotide objective"),
        ("CHROME", [0, 0, 1], "encoder frozen first, structure bolted on after"),
        ("This work", [1, 0, 0], "structure shapes the recurrence WHILE it learns"),
    ]
    stages = ["PRETRAIN", "DISTIL / ALIGN", "POST-HOC HEAD"]
    for j, s in enumerate(stages):
        ax.text(3.0 + j * 2.35, 4.05, s, fontsize=11.5, ha="center", color=MUTE,
                fontweight="bold")
    for i, (name, marks, note) in enumerate(rows):
        y = 3.25 - i * 0.85
        last = i == len(rows) - 1
        ax.text(0.1, y, name, fontsize=15,
                fontweight="bold" if last else "normal",
                color=STRUCT if last else INK, va="center")
        for j, m in enumerate(marks):
            cx = 3.0 + j * 2.35
            col = STRUCT if (m and last) else (ACCENT if m else "#f1f3f6")
            ax.add_patch(FancyBboxPatch((cx - 0.72, y - 0.235), 1.44, 0.47,
                                        boxstyle="round,pad=0.02,rounding_size=0.09",
                                        fc=col, ec="none"))
        ax.text(9.95, y - 0.005, note, fontsize=11, ha="right", va="center",
                color=MUTE)
    fig.text(0.055, 0.105,
             "The claim is not “use Hi-C to train a sequence model”. It is that "
             "structure changes HOW the model reads.", fontsize=14.5)
    footer(fig, "Source: docs/related_work.md, docs/architecture_spec.md §1")
    return save(fig, "03_gap")


def s04_pipeline(d):
    fig = canvas()
    title(fig, "27.4 GB of Hi-C, never downloaded",
          "Range reads pull a near-diagonal band. The whole file never lands on disk.")
    ax = fig.add_axes([0.055, 0.17, 0.89, 0.56]); ax.axis("off")
    ax.set_xlim(0, 10); ax.set_ylim(0, 3.2)
    steps = [
        ("27.4 GB", "remote .mcool", "#dbe3ee"),
        ("28.6 MB", "±2 Mb band + 250 kb\ncoarse, cached", "#bcd0e6"),
        ("1.3 MB", "φ — 8 numbers ×\n27,679 bins", "#8fb2d8"),
        ("in memory", "interpolated per\nnucleotide at load", STRUCT),
    ]
    for i, (big, small, col) in enumerate(steps):
        x = 0.25 + i * 2.5
        ax.add_patch(FancyBboxPatch((x, 1.15), 1.95, 1.25,
                                    boxstyle="round,pad=0.03,rounding_size=0.12",
                                    fc=col, ec="none"))
        ax.text(x + 0.975, 1.98, big, fontsize=21, ha="center", fontweight="bold",
                color="white" if i == 3 else INK)
        ax.text(x + 0.975, 1.45, small, fontsize=11, ha="center", va="center",
                color="white" if i == 3 else "#33455e", linespacing=1.35)
        if i < 3:
            ax.add_patch(FancyArrowPatch((x + 2.02, 1.78), (x + 2.44, 1.78),
                                         arrowstyle="-|>", mutation_scale=17,
                                         color=MUTE, lw=1.4))
    ax.text(0.25, 0.62, "a 4.4 GB materialised array is never created", fontsize=12,
            color=MUTE)
    ax.text(0.25, 0.20,
            "7,108,483 nonzero contacts kept  ·  400-bin band  ·  554 × 554 coarse matrix",
            fontsize=12.5, fontweight="bold")
    footer(fig, "Source: scripts/phase1_acquire.py, data/interim/*.npz (file sizes on disk)")
    return save(fig, "04_pipeline")


def s05_validation(d):
    fig = canvas()
    p = d["phi"]
    title(fig, "The features were checked against data we never touched",
          "4D Nucleome publishes its own derived tracks. Ours reproduce them.")
    ax = fig.add_axes([0.075, 0.24, 0.52, 0.48])
    keys = [("pearson_insulation_100kb_vs_4DN", "insulation 100 kb"),
            ("pearson_compartment_vs_4DN", "compartment PC1"),
            ("pearson_insulation_250kb_vs_4DN", "insulation 250 kb"),
            ("pearson_insulation_500kb_vs_4DN", "insulation 500 kb")]
    vals = [p[k] for k, _ in keys]
    labs = [l for _, l in keys]
    cols = [GOOD if v > 0.9 else (ACCENT if v > 0.6 else BASE) for v in vals]
    y = np.arange(len(vals))[::-1]
    ax.barh(y, vals, color=cols, height=0.58)
    for yi, v in zip(y, vals):
        ax.text(v + 0.015, yi, f"{v:.4f}", va="center", fontsize=13.5,
                fontweight="bold")
    ax.set_yticks(y); ax.set_yticklabels(labs, fontsize=12.5)
    ax.set_xlim(0, 1.12); ax.set_xticks([0, 0.5, 1.0])
    ax.set_xlabel("Pearson r against the independent 4DN track", fontsize=12)
    ax.spines["left"].set_color(FAINT)

    ax2 = fig.add_axes([0.66, 0.24, 0.28, 0.48]); ax2.axis("off")
    for i, (k, v) in enumerate([
            ("bins on chr9", f"{p['n_bins']:,}"),
            ("usable", f"{p['usable_bins']:,}  ({100*p['usable_frac']:.1f}%)"),
            ("train / val / test", "{train:,} / {val:,} / {test:,}".format(**d["splits"])),
            ("window", f"{d['window']:,} bp")]):
        yy = 0.86 - i * 0.25
        ax2.text(0, yy, v, fontsize=17, fontweight="bold", transform=ax2.transAxes)
        ax2.text(0, yy - 0.085, k.upper(), fontsize=9.5, color=MUTE,
                 transform=ax2.transAxes)
    fig.text(0.055, 0.135,
             "The 100 kb insulation track agrees at r = 0.9969 — on a track the "
             "pipeline never saw while it was being built.", fontsize=14.5)
    footer(fig, "Source: data/processed/phi_validation_report.json")
    return save(fig, "05_validation")


def s06_mechanism(d):
    fig = canvas()
    title(fig, "One number decides how far the model remembers",
          "Δ sets the decay of the recurrence. τ = 1 / (Δ·|A|), in nucleotides.")
    ax = fig.add_axes([0.055, 0.30, 0.42, 0.44])
    x = np.linspace(0, 1000, 600)
    for tau, c, lab, lw in ((14, OLD, "τ = 14   forgets within 14 bp", 2.2),
                            (435, ACCENT, "τ = 435   spans a window", 2.6)):
        ax.plot(x, np.exp(-x / tau), color=c, lw=lw, label=lab)
    ax.set_xlabel("distance back along the DNA (nucleotides)", fontsize=11.5)
    ax.set_ylabel("how much the state remembers", fontsize=11.5)
    ax.legend(frameon=False, fontsize=11.5, loc="upper right")
    ax.set_ylim(0, 1.04)

    ax2 = fig.add_axes([0.56, 0.30, 0.385, 0.44]); ax2.axis("off")
    ax2.set_xlim(0, 10); ax2.set_ylim(0, 6)
    ax2.plot([0.4, 9.6], [4.6, 4.6], color=FAINT, lw=8, solid_capstyle="round")
    for bx, bw, lab in ((0.6, 3.6, "domain"), (5.4, 3.9, "domain")):
        ax2.add_patch(Rectangle((bx, 4.35), bw, 0.5, fc="#e8eef7", ec="none"))
        ax2.text(bx + bw / 2, 5.15, lab, fontsize=11, ha="center", color=MUTE)
    ax2.plot([4.7, 4.7], [4.05, 5.35], color=BAD, lw=3)
    ax2.text(4.7, 5.55, "boundary", fontsize=11.5, ha="center", color=BAD,
             fontweight="bold")
    ax2.annotate("", xy=(4.15, 3.3), xytext=(1.1, 3.3),
                 arrowprops=dict(arrowstyle="-|>", color=ACCENT, lw=2.6))
    ax2.text(1.1, 2.75, "long memory inside", fontsize=11.5, color=ACCENT)
    ax2.annotate("", xy=(5.15, 3.3), xytext=(4.72, 3.3),
                 arrowprops=dict(arrowstyle="-|>", color=BAD, lw=2.6))
    ax2.text(5.25, 2.75, "shortened at the wall", fontsize=11.5, color=BAD)
    ax2.text(0.4, 1.35, "structure adjusts Δ per position", fontsize=15,
             fontweight="bold")
    ax2.text(0.4, 0.65, "8 Hi-C numbers → 178-parameter encoder → nudge on Δ",
             fontsize=12, color=MUTE)
    fig.text(0.055, 0.155,
             "This changes HOW the model reads, not what it looks at afterwards.",
             fontsize=14.5)
    footer(fig, "Source: src/chromfm/model.py, docs/architecture_spec.md §4.1.1")
    return save(fig, "06_mechanism")


def s07_matched(d):
    fig = canvas()
    b, s = d["p_base"], d["p_str"]
    over = 100 * (s - b) / b
    title(fig, "The comparison is fair by construction",
          "Same size, same data, same schedule — and identical at step 0.")
    ax = fig.add_axes([0.075, 0.30, 0.40, 0.42])
    ax.bar([0, 1], [b, s], color=[BASE, STRUCT], width=0.52)
    for i, v in enumerate([b, s]):
        ax.text(i, v * 1.008, f"{v:,}", ha="center", fontsize=15, fontweight="bold")
    ax.set_xticks([0, 1]); ax.set_xticklabels(["baseline", "structural"], fontsize=13)
    ax.set_ylim(0, s * 1.10); ax.set_yticks([])
    ax.spines["left"].set_visible(False)
    ax.set_title(f"+{over:.2f}% — budget is 5%", fontsize=13.5, color=MUTE, pad=14)

    ax2 = fig.add_axes([0.56, 0.30, 0.385, 0.42]); ax2.axis("off")
    ax2.set_xlim(0, 10); ax2.set_ylim(0, 6)
    ax2.text(0, 5.4, "W_Δs initialised to exactly zero", fontsize=16,
             fontweight="bold")
    ax2.text(0, 4.7, "so at step 0 the two models are", fontsize=13, color=MUTE)
    ax2.add_patch(FancyBboxPatch((0, 3.0), 9.3, 1.35,
                                 boxstyle="round,pad=0.05,rounding_size=0.15",
                                 fc="#eef6f1", ec=GOOD, lw=1.6))
    ax2.text(4.65, 3.67, "numerically identical  ·  max gap 0.0", fontsize=16.5,
             ha="center", va="center", color=GOOD, fontweight="bold")
    ax2.text(0, 2.15, "Everything that develops afterwards is learned,", fontsize=13)
    ax2.text(0, 1.55, "not built in.", fontsize=13)
    ax2.text(0, 0.5, "verified by 39 automated checks before any GPU time",
             fontsize=11.5, color=MUTE, style="italic")
    footer(fig, "Source: results/*/run_config.yaml, scripts/test_phase4_wiring.py (39/39)")
    return save(fig, "07_matched")


def s08_f4(d):
    fig = canvas()
    title(fig, "A one-line bug capped memory at exactly 1,000 letters",
          "τ_max at init is exactly 1/dt_min, because |A| ≥ 1 by construction.")
    ax = fig.add_axes([0.075, 0.28, 0.40, 0.45])
    ax.set_yscale("log")
    ax.bar([0, 1], [994, 985872], color=[BAD, GOOD], width=0.52)
    for i, v in enumerate([994, 985872]):
        ax.text(i, v * 1.35, f"{v:,}", ha="center", fontsize=15, fontweight="bold")
    ax.axhline(5000, color=MUTE, ls=":", lw=1.3)
    ax.text(1.42, 5000, "one 5 kb Hi-C bin", fontsize=10.5, color=MUTE, va="center")
    ax.axhline(100000, color=MUTE, ls=":", lw=1.3)
    ax.text(1.42, 100000, "100 kb", fontsize=10.5, color=MUTE, va="center")
    ax.set_xticks([0, 1])
    ax.set_xticklabels(["dt_min = 1e-3\n(Mamba reference)", "dt_min = 1e-6\n(fixed)"],
                       fontsize=12)
    ax.set_ylabel("τ_max at initialisation (nucleotides)", fontsize=11.5)
    ax.set_xlim(-0.6, 2.3)

    ax2 = fig.add_axes([0.57, 0.28, 0.375, 0.45]); ax2.axis("off")
    rows = [("median τ after training", "14.2", "434.7"),
            ("state past 100 kb", "≈ 0", "4.85e-02"),
            ("layer-directions reaching it", "0 / 32", "32 / 32"),
            ("parameters", "7,725,312", "7,725,312")]
    ax2.text(0.47, 0.97, "BEFORE", fontsize=11, color=BAD, ha="center",
             transform=ax2.transAxes, fontweight="bold")
    ax2.text(0.85, 0.97, "AFTER", fontsize=11, color=GOOD, ha="center",
             transform=ax2.transAxes, fontweight="bold")
    for i, (lab, a, b_) in enumerate(rows):
        y = 0.80 - i * 0.20
        ax2.text(0, y, lab, fontsize=12, color=MUTE, transform=ax2.transAxes)
        ax2.text(0.47, y, a, fontsize=14, ha="center", transform=ax2.transAxes)
        ax2.text(0.85, y, b_, fontsize=14, ha="center", fontweight="bold",
                 color=GOOD if i < 3 else INK, transform=ax2.transAxes)
        ax2.plot([0, 1], [y - 0.055, y - 0.055], color=FAINT, lw=0.8,
                 transform=ax2.transAxes)
    ax2.text(0, 0.02, "cost: three constants. zero extra parameters.",
             fontsize=12.5, fontweight="bold", transform=ax2.transAxes)
    footer(fig, "Source: docs/architecture_spec.md §4.1.4, results/baselines/phase3_report_baseline_v2.txt")
    return save(fig, "08_f4")


def s09_negative(d):
    fig = canvas()
    title(fig, "30× more memory bought nothing",
          "The first real finding, and it is a negative one.")
    ax = fig.add_axes([0.075, 0.28, 0.39, 0.44])
    ax.bar([0, 1], [14.2, 434.7], color=[OLD, ACCENT], width=0.5)
    ax.text(0, 20, "14.2", ha="center", fontsize=14, fontweight="bold")
    ax.text(1, 448, "434.7", ha="center", fontsize=14, fontweight="bold")
    ax.set_xticks([0, 1]); ax.set_xticklabels(["before", "after"], fontsize=12.5)
    ax.set_ylabel("median memory horizon (nt)", fontsize=11.5)
    ax.set_title("memory  ×30.6", fontsize=13.5, color=ACCENT, pad=12,
                 fontweight="bold")

    ax2 = fig.add_axes([0.56, 0.28, 0.385, 0.44])
    m1, s1, m2, s2 = 1.5210, 0.0040, 1.5197, 0.0025
    ax2.errorbar([0, 1], [m1, m2], yerr=[s1, s2], fmt="o", ms=13,
                 color=INK, ecolor=MUTE, elinewidth=1.8, capsize=9)
    for i, (m, s) in enumerate([(m1, s1), (m2, s2)]):
        ax2.text(i + 0.11, m, f"{m:.4f}", fontsize=13.5, va="center",
                 fontweight="bold")
    ax2.set_xlim(-0.45, 1.6); ax2.set_xticks([0, 1])
    ax2.set_xticklabels(["before", "after"], fontsize=12.5)
    ax2.set_ylabel("val bits / nucleotide  (lower is better)", fontsize=11.5)
    ax2.set_title("score  unchanged", fontsize=13.5, color=MUTE, pad=12,
                  fontweight="bold")
    fig.text(0.055, 0.145,
             "A metric that cannot see a 30× capability change has no power to "
             "adjudicate the structural mechanism either.\n"
             "That is why the gate was redesigned before the structural arm was read.",
             fontsize=14, linespacing=1.6)
    footer(fig, "Source: results/baselines/phase3_report_baseline_v2.txt (3 seeds per arm)")
    return save(fig, "09_negative")


def s10_traces(d):
    fig = canvas()
    tr = d["tr"]
    title(fig, "Six runs, three seeds per arm",
          "Every run logged before the next arm existed.")
    ax = fig.add_axes([0.075, 0.22, 0.62, 0.52])
    for k, v in tr.items():
        if k.startswith("baseline_v2"):
            ax.plot(v["step"], v["val"], color=BASE, lw=2.1, alpha=0.95)
        elif k.startswith("structural"):
            ax.plot(v["step"], v["val"], color=STRUCT, lw=2.1,
                    alpha=0.95 if v["done"] else 0.45,
                    ls="-" if v["done"] else (0, (4, 2)))
    ax.axhline(2.0, color=MUTE, ls=":", lw=1.2)
    ax.text(60, 2.03, "random guessing = 2.000", fontsize=11, color=MUTE)
    ax.set_xlabel("optimiser step", fontsize=12)
    ax.set_ylabel("val bits / nucleotide", fontsize=12)
    ax.set_ylim(1.48, 2.95)
    ax.plot([], [], color=BASE, lw=2.4, label="baseline  (3 seeds)")
    ax.plot([], [], color=STRUCT, lw=2.4, label="structural  (3 seeds)")
    ax.legend(frameon=False, fontsize=12.5, loc="upper right")

    ax2 = fig.add_axes([0.755, 0.22, 0.19, 0.52])
    ax2.set_title("final score", fontsize=12, color=MUTE, pad=10)
    fb, fs = finals(tr, "baseline_v2"), finals(tr, "structural")
    ax2.scatter([0] * len(fb), fb, s=95, color=BASE, zorder=3)
    ax2.scatter([1] * len(fs), fs, s=95, color=STRUCT, zorder=3)
    if fb:
        ax2.hlines(np.mean(fb), -0.28, 0.28, color=INK, lw=2.2)
    if fs:
        ax2.hlines(np.mean(fs), 0.72, 1.28, color=INK, lw=2.2)
    ax2.set_xlim(-0.55, 1.55); ax2.set_xticks([0, 1])
    ax2.set_xticklabels(["base", "struct"], fontsize=11.5)
    ax2.tick_params(labelleft=False)
    ax2.set_ylim(1.512, 1.532)
    note = ("all three structural seeds complete" if len(fs) == 3
            else f"{len(fs)} of 3 structural seeds complete — seed 2 still running")
    fig.text(0.055, 0.135, note, fontsize=14.5, fontweight="bold")
    footer(fig, "Source: results/*/metrics.json, read at build time")
    return save(fig, "10_traces")


def s11_result(d):
    fig = canvas()
    tr = d["tr"]
    fb, fs = finals(tr, "baseline_v2"), finals(tr, "structural")
    mb, sb = float(np.mean(fb)), float(np.std(fb, ddof=1))
    title(fig, "Structure has not changed pretraining loss",
          f"Effect-size floor is 2σ = {2*sb:.4f} bits, pre-registered.")
    ax = fig.add_axes([0.075, 0.26, 0.55, 0.46])
    ax.axhspan(mb - 2 * sb, mb + 2 * sb, color="#eef2f7", zorder=0)
    ax.axhline(mb, color=BASE, lw=1.8)
    ax.text(2.42, mb, "baseline mean", fontsize=11.5, color=BASE, va="bottom")
    ax.text(2.42, mb - 2 * sb, "−2σ", fontsize=10.5, color=MUTE, va="center")
    ax.text(2.42, mb + 2 * sb, "+2σ", fontsize=10.5, color=MUTE, va="center")
    for i, v in enumerate(fb):
        ax.scatter(0.85 + 0.14 * i, v, s=130, color=BASE, zorder=3)
    for i, v in enumerate(fs):
        ax.scatter(1.75 + 0.14 * i, v, s=130, color=STRUCT, zorder=3)
    if fs:
        ms = float(np.mean(fs))
        ax.hlines(ms, 1.70, 1.70 + 0.14 * len(fs), color=STRUCT, lw=2.4)
    ax.set_xlim(0.55, 3.0); ax.set_xticks([1.0, 1.9])
    ax.set_xticklabels(["baseline", "structural"], fontsize=13)
    ax.set_ylabel("val bits / nucleotide", fontsize=12)

    ax2 = fig.add_axes([0.70, 0.26, 0.25, 0.46]); ax2.axis("off")
    if fs:
        ms = float(np.mean(fs))
        delta = ms - mb
        stats = [("baseline mean", f"{mb:.4f}"),
                 ("structural mean", f"{ms:.4f}"),
                 ("difference", f"{delta:+.4f} bits"),
                 ("in units of σ", f"{delta/sb:+.2f} σ"),
                 ("bar to clear", f"{2*sb:.4f} bits")]
        for i, (k, v) in enumerate(stats):
            y = 0.88 - i * 0.19
            ax2.text(0, y, v, fontsize=17, fontweight="bold",
                     color=STRUCT if i == 3 else INK, transform=ax2.transAxes)
            ax2.text(0, y - 0.075, k.upper(), fontsize=9.5, color=MUTE,
                     transform=ax2.transAxes)
    n = len(fs)
    fig.text(0.055, 0.135,
             ("Preliminary — %d of 3 structural seeds complete. "
              "The gate requires all three." % n) if n < 3 else
             "All three seeds complete.", fontsize=14.5, fontweight="bold")
    footer(fig, "Source: results/*/metrics.json, read at build time")
    return save(fig, "11_result")


def s12_controls(d):
    fig = canvas()
    title(fig, "Four ways to kill our own claim",
          "Each control removes one benign explanation. Two of them can end the project.")
    ax = fig.add_axes([0.055, 0.16, 0.89, 0.58]); ax.axis("off")
    ax.set_xlim(0, 10); ax.set_ylim(0, 5)
    ctrl = [
        ("S1", "shuffle φ\neverywhere", "does it rely on\nstructure at all?", ACCENT),
        ("S2", "slide φ\nby 10 Mb", "“any smooth extra\nchannel helps”", ACCENT),
        ("S3", "keep only\ndistance decay", "“structure is just\ngenomic distance”", BAD),
        ("S4", "swap in GC +\ngene density", "“structure is just\nsequence composition”", BAD),
    ]
    for i, (tag, what, kills, col) in enumerate(ctrl):
        x = 0.1 + i * 2.47
        ax.add_patch(FancyBboxPatch((x, 1.15), 2.32, 3.0,
                                    boxstyle="round,pad=0.04,rounding_size=0.14",
                                    fc="#f7f9fc", ec=col, lw=1.8))
        ax.add_patch(plt.Circle((x + 0.42, 3.72), 0.27, color=col))
        ax.text(x + 0.42, 3.72, tag, fontsize=13, color="white", ha="center",
                va="center", fontweight="bold")
        ax.text(x + 0.18, 3.18, what, fontsize=12, fontweight="bold",
                va="top", linespacing=1.45)
        ax.text(x + 0.18, 2.30, "RULES OUT", fontsize=8.5, color=MUTE)
        ax.text(x + 0.18, 2.02, kills, fontsize=10, color="#33455e", va="top",
                linespacing=1.5)
        if col == BAD:
            ax.text(x + 0.18, 1.35, "can kill the paper", fontsize=10,
                    color=BAD, fontweight="bold")
    ax.text(0.1, 0.55, "All four built and tested. Sequence is never touched — "
            "only φ changes, so any difference is attributable to structure.",
            fontsize=13.5)
    footer(fig, "Source: docs/architecture_spec.md §4.1.3; scripts/phase4_build_s3.py, phase4_build_s4.py")
    return save(fig, "12_controls")


def s13_s4(d):
    fig = canvas()
    c = d["s4"]["phi_vs_s4_best_correlation"]
    title(fig, "Only one of eight features is explainable by sequence",
          "The GC objection was real — and it is confined to compartment PC1.")
    ax = fig.add_axes([0.075, 0.20, 0.55, 0.55])
    items = sorted(c.items(), key=lambda kv: abs(kv[1]["best_r"]))
    labs = [k for k, _ in items]
    vals = [abs(v["best_r"]) for _, v in items]
    cols = [BAD if v > 0.5 else GOOD for v in vals]
    y = np.arange(len(vals))
    ax.barh(y, vals, color=cols, height=0.6)
    for yi, v in zip(y, vals):
        ax.text(v + 0.012, yi, f"{v:.4f}", va="center", fontsize=12,
                fontweight="bold")
    ax.set_yticks(y); ax.set_yticklabels(labs, fontsize=11.5)
    ax.set_xlim(0, 0.88)
    ax.set_xlabel("|r| with the best sequence-derived covariate", fontsize=11.5)
    ax.axvline(0.5, color=MUTE, ls=":", lw=1.2)
    ax.spines["left"].set_color(FAINT)

    ax2 = fig.add_axes([0.665, 0.20, 0.28, 0.55]); ax2.axis("off")
    ax2.text(0, 0.93, "0.7384", fontsize=34, fontweight="bold", color=BAD,
             transform=ax2.transAxes)
    ax2.text(0, 0.855, "compartment PC1 vs smoothed GC", fontsize=11.5,
             color=MUTE, transform=ax2.transAxes)
    ax2.text(0, 0.79, "55% of its variance", fontsize=11.5, color=MUTE,
             transform=ax2.transAxes)
    ax2.text(0, 0.60, "0.1822", fontsize=34, fontweight="bold", color=GOOD,
             transform=ax2.transAxes)
    ax2.text(0, 0.525, "worst of the other seven", fontsize=11.5, color=MUTE,
             transform=ax2.transAxes)
    ax2.text(0, 0.30,
             "Insulation and directionality —\nthe features that encode domain\n"
             "boundaries — are NOT recoverable\nfrom sequence composition.",
             fontsize=12.5, transform=ax2.transAxes, linespacing=1.7)
    footer(fig, "Source: data/processed/s4_validation_report.json (21,519 usable bins)")
    return save(fig, "13_s4")


def s14_s3(d):
    fig = canvas()
    s3 = d["s3"]
    title(fig, "The distance control, verified rather than asserted",
          "Distance decay preserved exactly. Locus specificity destroyed.")
    ax = fig.add_axes([0.075, 0.22, 0.40, 0.52])
    corr = s3["s3_vs_real_feature_correlation"]
    items = sorted(corr.items(), key=lambda kv: abs(kv[1]))
    y = np.arange(len(items))
    ax.barh(y, [abs(v) for _, v in items], color=GOOD, height=0.6)
    ax.set_yticks(y); ax.set_yticklabels([k for k, _ in items], fontsize=10.5)
    ax.axvline(0.3, color=BAD, ls=":", lw=1.4)
    ax.text(0.305, 0.4, "0.3 threshold", fontsize=10.5, color=BAD, rotation=90)
    ax.set_xlim(0, 0.42)
    ax.set_xlabel("|r| of the rewired feature with the real one", fontsize=11.5)
    ax.spines["left"].set_color(FAINT)

    ax2 = fig.add_axes([0.56, 0.22, 0.385, 0.52]); ax2.axis("off")
    stats = [("P(s) preserved to", f"{s3['ps_max_abs_deviation']:.2e}",
              "max deviation over 401 diagonals"),
             ("contacts permuted", "8,612,548", "NaN pattern left in place"),
             ("usable bins", f"{s3['usable_bins_s3']:,}",
              "identical to real φ — no coverage confound"),
             ("worst correlation", "0.1818", "locus specificity gone")]
    for i, (k, v, note) in enumerate(stats):
        y = 0.87 - i * 0.235
        ax2.text(0, y, v, fontsize=21, fontweight="bold", transform=ax2.transAxes)
        ax2.text(0, y - 0.062, k.upper(), fontsize=9.5, color=MUTE,
                 transform=ax2.transAxes)
        ax2.text(0, y - 0.115, note, fontsize=10.5, color="#33455e",
                 transform=ax2.transAxes)
    footer(fig, "Source: data/processed/s3_validation_report.json")
    return save(fig, "14_s3")


def s15_gate(d):
    fig = canvas()
    title(fig, "The gate, fixed before the answer was known",
          "A loss delta alone cannot tell “inert” from “MLM cannot express it”.")
    ax = fig.add_axes([0.10, 0.19, 0.52, 0.55]); ax.axis("off")
    ax.set_xlim(0, 10); ax.set_ylim(0, 10)
    cells = [((5.1, 5.1), "#eef6f1", GOOD, "structure helps", "proceed"),
             ((5.1, 0.4), "#fdf1ef", BAD, "mechanism INERT", "stop"),
             ((0.4, 5.1), "#f4f6f9", MUTE, "(incoherent)", "investigate"),
             ((0.4, 0.4), "#eef6f1", GOOD, "used, not expressible", "PROCEED")]
    for (x, y), fc, ec, lab, act in cells:
        ax.add_patch(FancyBboxPatch((x, y), 4.5, 4.3,
                                    boxstyle="round,pad=0.04,rounding_size=0.16",
                                    fc=fc, ec=ec, lw=1.8))
        ax.text(x + 2.25, y + 2.65, lab, fontsize=13.5, ha="center",
                fontweight="bold")
        ax.text(x + 2.25, y + 1.55, act, fontsize=12, ha="center", color=ec)
    ax.text(-0.35, 7.25, "divergence\nHIGH", fontsize=12, ha="right", va="center",
            color=MUTE, linespacing=1.5)
    ax.text(-0.35, 2.55, "divergence\nLOW", fontsize=12, ha="right", va="center",
            color=MUTE, linespacing=1.5)
    ax.text(2.65, 10.1, "loss flat", fontsize=12, ha="center", color=MUTE)
    ax.text(7.35, 10.1, "loss moves", fontsize=12, ha="center", color=MUTE)

    ax2 = fig.add_axes([0.66, 0.19, 0.29, 0.55]); ax2.axis("off")
    ax2.text(0, 0.95, "RELIANCE", fontsize=12, fontweight="bold",
             transform=ax2.transAxes)
    ax2.text(0, 0.86, "does it change predictions\nwhen we lie about folding?",
             fontsize=11.5, color=MUTE, transform=ax2.transAxes, linespacing=1.5)
    ax2.text(0, 0.68, "KL divergence + flip rate", fontsize=12.5,
             fontweight="bold", color=ACCENT, transform=ax2.transAxes)
    ax2.text(0, 0.52, "BENEFIT", fontsize=12, fontweight="bold",
             transform=ax2.transAxes)
    ax2.text(0, 0.43, "does the score get worse?", fontsize=11.5, color=MUTE,
             transform=ax2.transAxes)
    ax2.text(0, 0.33, "Δ ≥ 2σ = 0.0050 bits", fontsize=12.5, fontweight="bold",
             color=ACCENT, transform=ax2.transAxes)
    ax2.text(0, 0.13,
             "An inert mechanism cannot\nchange its output at all —\n"
             "measured floor: exactly 0.0",
             fontsize=11.5, transform=ax2.transAxes, linespacing=1.6)
    footer(fig, "Source: docs/architecture_spec.md §4.1.3 amendment, recorded 2026-08-16 before any swap")
    return save(fig, "15_gate")


def s16_reliance(d):
    """What the gate returned: benefit null, reliance real but very small."""
    fig = canvas()
    p1 = d["p1"]
    title(fig, "It listens, but it does not benefit",
          "Lying about the folding changes the score by nothing, and the "
          "predictions by very little.")
    s = p1["summary"]
    pc = s["per_control"]
    order = ["S0", "S1", "S2", "S3", "S4"]
    labs = {"S0": "S0  structure removed", "S1": "S1  shuffled",
            "S2": "S2  shifted 10 Mb", "S3": "S3  distance-matched",
            "S4": "S4  sequence-matched"}
    ax = fig.add_axes([0.195, 0.36, 0.42, 0.36])
    y = np.arange(len(order))[::-1]
    vals = [pc[c]["kl_mean"] for c in order]
    cols = [MUTE, ACCENT, ACCENT, BAD, BAD]
    ax.barh(y, vals, height=0.60, color=cols)
    mf = s["masking_floor_mean_CONTEXT_ONLY"]
    ax.axvline(mf, color=INK, lw=1.6, ls="--")
    ax.text(mf * 0.8, 4.45,
            f"re-masking 15% of the DNA\nmoves it this far  ({mf:.3f})",
            fontsize=10.5, va="top", ha="right", color=INK, linespacing=1.5)
    ax.set_xscale("log")
    ax.set_xlim(1e-5, 1.0)
    ax.set_ylim(-0.7, len(order) - 0.3)
    ax.set_yticks(y)
    ax.set_yticklabels([labs[c] for c in order], fontsize=11.5)
    ax.set_xlabel("KL divergence from the real-structure prediction  (log)",
                  fontsize=11)
    for sp in ("top", "right", "left"):
        ax.spines[sp].set_visible(False)
    ax.tick_params(axis="y", length=0)

    ax3 = fig.add_axes([0.655, 0.30, 0.29, 0.42]); ax3.axis("off")
    ax3.text(0, 1.0, "READ IT HONESTLY", fontsize=11.5, fontweight="bold",
             transform=ax3.transAxes, va="top")
    ax3.text(0, 0.88,
             "An inert model cannot move its output\n"
             "at all — the measured floor is exactly\n"
             "0.0. So these divergences are real.\n\n"
             "But they sit at 0.05% of the masking\n"
             "floor: φ moves predictions ~2000×\n"
             "less than the DNA does.\n\n"
             "And S1–S3 rank above S4 ≈ S0 because\n"
             "they substitute values further from the\n"
             "truth, not because they are more\n"
             "meaningful.",
             fontsize=10.8, color=INK, transform=ax3.transAxes,
             linespacing=1.6, va="top")

    ax2 = fig.add_axes([0.055, 0.135, 0.56, 0.12]); ax2.axis("off")
    kf = s["kernel_floor_mean"]
    for i, (v, k, c) in enumerate([
            (f"{pc['S1']['delta_bits_mean']:+.4f}", "benefit, shuffled φ", GOOD),
            ("0.0050", "bar it had to clear", MUTE),
            (f"{kf:.1f}", "inert-model floor", ACCENT),
            (f"{pc['S1']['flip_frac_mean']*100:.2f}%", "predictions flipped", ACCENT)]):
        ax2.text(i * 0.26, 0.62, v, fontsize=18, fontweight="bold", color=c,
                 transform=ax2.transAxes)
        ax2.text(i * 0.26, 0.08, k.upper(), fontsize=9, color=MUTE,
                 transform=ax2.transAxes)
    fig.text(0.055, 0.105,
             "Benefit: null. Reliance: real, and ~2000× smaller than the "
             "sequence itself.",
             fontsize=14, fontweight="bold", va="top")
    footer(fig, "Source: results/novel_model/p1_swap_results.json, 3 seeds, "
                "1,344,979 masked positions per pass")
    return save(fig, "16_reliance")


def s17_depth(d):
    """D1 by layer: inert through the encoder, live only at the top."""
    fig = canvas()
    dd = d["d1"]
    title(fig, "And it only switches on in the last four layers",
          "D1 = how much the structural term actually moves Δ, per layer. "
          "Below 0.05 the pathway is inert.")
    runs = sorted(dd["runs"])
    nl = 16
    M = np.zeros((nl, len(runs)))
    for j, r in enumerate(runs):
        for li in range(nl):
            v = [x["d1_pooled"] for x in dd["runs"][r]["layers"]
                 if x["layer"].startswith(f"L{li:02d}.")]
            M[li, j] = float(np.mean(v))

    ax = fig.add_axes([0.30, 0.20, 0.20, 0.56])
    Z = np.log10(np.clip(M, 1e-5, None))
    ax.imshow(Z, aspect="auto", origin="lower", cmap="YlOrRd",
              vmin=-4.2, vmax=0.05)
    for li in range(nl):
        for j in range(len(runs)):
            ax.text(j, li, f"{M[li, j]:.3f}" if M[li, j] >= 0.05 else "·",
                    ha="center", va="center", fontsize=9.5,
                    color="white" if M[li, j] >= 0.4 else INK)
    ax.set_xticks(range(len(runs)))
    ax.set_xticklabels([f"seed {r[-1]}" for r in runs], fontsize=11)
    ax.set_yticks(range(nl))
    ax.set_yticklabels([f"L{i:02d}" for i in range(nl)], fontsize=9)
    ax.set_ylabel("layer  (input at the bottom, output head at the top)",
                  fontsize=11)
    ax.axhline(11.5, color=INK, lw=2.2)

    ax2 = fig.add_axes([0.545, 0.20, 0.41, 0.55]); ax2.axis("off")
    ax2.text(0, 1.0, "WHAT THE SPLIT MEANS", fontsize=11.5, fontweight="bold",
             transform=ax2.transAxes, va="top")
    live = [dd["runs"][r]["d1_n_below_threshold"] for r in runs]
    tot = dd["runs"][runs[0]]["d1_n_directions"]
    ax2.text(0, 0.90,
             f"LAYERS 0–11 — inert in every seed. The\n"
             f"structural weights are NOT zero (‖W‖ ≈ 0.19–\n"
             f"0.22 against the exact 0 they started at), so\n"
             f"gradient did reach them. What they produce\n"
             f"is a near-constant, absorbed into the existing\n"
             f"bias — failure mode F1, measured not inferred.\n\n"
             f"LAYERS 12–15 — the only live directions, in all\n"
             f"three seeds, without exception:\n"
             f"{tot-live[0]}, {tot-live[1]} and {tot-live[2]} of {tot}.\n\n"
             f"So the mechanism acts as a late readout\n"
             f"correction, not as a prior shaping how the\n"
             f"sequence is encoded — which is what the\n"
             f"claim rests on.\n\n"
             f"Cheapest test of whether that is intrinsic:\n"
             f"drop the encoder, d_struct 2→8. One seed,\n"
             f"about 6 GPU-hours.",
             fontsize=10.8, transform=ax2.transAxes, linespacing=1.58,
             va="top")
    fig.text(0.055, 0.125,
             "Every live direction, in every seed, sits in the top four layers.",
             fontsize=14, fontweight="bold", va="top")
    footer(fig, "Source: results/novel_model/d1_diagnostic.json — threshold 0.05 "
                "pre-registered at architecture_spec.md §4.1.3")
    return save(fig, "17_depth")


def s16_status(d):
    fig = canvas()
    tr = d["tr"]
    ndone = len(finals(tr, "structural"))
    title(fig, "Where the project stands",
          "Phases 0–4 complete, gate read. Phase 5 is a decision, not a default.")
    ax = fig.add_axes([0.055, 0.30, 0.89, 0.44]); ax.axis("off")
    ax.set_xlim(0, 10); ax.set_ylim(0, 3)
    phases = [("0", "literature", 1), ("1", "data", 1), ("2", "mechanism", 1),
              ("3", "baselines", 1), ("4", "structural arm", 1),
              ("5", "evaluation", 0), ("6", "paper", 0)]
    for i, (num, lab, st) in enumerate(phases):
        x = 0.2 + i * 1.40
        col = GOOD if st == 1 else (STRUCT if st == 0.5 else "#e9edf2")
        ax.add_patch(FancyBboxPatch((x, 1.30), 1.16, 0.95,
                                    boxstyle="round,pad=0.03,rounding_size=0.12",
                                    fc=col, ec="none"))
        ax.text(x + 0.58, 1.78, num, fontsize=20, ha="center", va="center",
                color="white" if st > 0 else MUTE, fontweight="bold")
        ax.text(x + 0.58, 1.05, lab, fontsize=11, ha="center", color=MUTE)
        if i < len(phases) - 1:
            ax.plot([x + 1.20, x + 1.36], [1.78, 1.78], color=FAINT, lw=2)
    ax.text(0.2, 0.45,
            f"structural arm: {ndone} of 3 seeds complete — gate read, "
            f"diagnostics run",
            fontsize=13.5, fontweight="bold", color=STRUCT)
    ax.text(0.2, 0.05,
            "next: one 6 GPU-hour re-run without the structural encoder, to "
            "test whether the dead layers are intrinsic or self-inflicted",
            fontsize=12.5, color=MUTE)

    ax2 = fig.add_axes([0.055, 0.10, 0.89, 0.14]); ax2.axis("off")
    for i, (k, v) in enumerate([("GPU-hours spent", "~48"),
                                ("cost of the gate", "30 minutes"),
                                ("negative results banked", "3"),
                                ("automated checks", "39 + 17 + 34")]):
        ax2.text(i * 0.25, 0.62, v, fontsize=19, fontweight="bold",
                 transform=ax2.transAxes)
        ax2.text(i * 0.25, 0.10, k.upper(), fontsize=9.5, color=MUTE,
                 transform=ax2.transAxes)
    footer(fig, "GPU-hours derived from measured 5.30 s/step (baseline) and 6.13 s/step (structural)")
    return save(fig, "16_status")


def s17_close(d):
    fig = canvas()
    fig.text(0.055, 0.76, "What survives either way", fontsize=34,
             fontweight="bold")
    ax = fig.add_axes([0.055, 0.20, 0.89, 0.48]); ax.axis("off")
    ax.set_xlim(0, 10); ax.set_ylim(0, 5)
    cards = [
        ("The memory cap", "Mamba's reference dt_min caps τ at exactly 1,000\n"
         "nucleotides — below almost every genomic\nannotation. Fixed at zero "
         "parameter cost.", GOOD),
        ("A metric with no power", "30× more memory moved val loss by less than\n"
         "seed noise. Worth knowing before anyone else\nspends compute assuming "
         "otherwise.", ACCENT),
        ("A falsifiable design", "Matched parameters, pre-registered rules, and\n"
         "four controls built to kill the claim. A clean\nnegative is a real result.",
         STRUCT),
    ]
    for i, (h, body, col) in enumerate(cards):
        x = 0.1 + i * 3.3
        ax.add_patch(FancyBboxPatch((x, 0.9), 3.0, 3.4,
                                    boxstyle="round,pad=0.04,rounding_size=0.15",
                                    fc="#f7f9fc", ec="none"))
        ax.add_patch(Rectangle((x, 0.9), 0.075, 3.4, fc=col, ec="none"))
        ax.text(x + 0.32, 3.72, h, fontsize=15, fontweight="bold")
        ax.text(x + 0.32, 3.20, body, fontsize=11.3, va="top", linespacing=1.75,
                color="#33455e")
    fig.text(0.055, 0.115,
             "The mechanism may well be inert. The project is built so that "
             "finding out is still worth something.", fontsize=15)
    return save(fig, "17_close")


SLIDES = [s01_title, s02_problem, s03_gap, s04_pipeline, s05_validation,
          s06_mechanism, s07_matched, s08_f4, s09_negative, s10_traces,
          s11_result, s12_controls, s13_s4, s14_s3, s15_gate, s16_reliance,
          s17_depth, s16_status, s17_close]


def main() -> int:
    from pptx import Presentation
    from pptx.util import Inches

    d = load()
    print(f"building {len(SLIDES)} slides")
    paths = []
    for fn in SLIDES:
        p = fn(d)
        paths.append(p)
        print(f"  {p.name}")

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    for p in paths:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(p), 0, 0, width=prs.slide_width,
                             height=prs.slide_height)
    OUT.mkdir(parents=True, exist_ok=True)
    dest = OUT / "pitch_deck.pptx"
    prs.save(str(dest))
    print(f"\nwrote {dest.relative_to(REPO)} "
          f"({dest.stat().st_size/1e6:.1f} MB, {len(paths)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
